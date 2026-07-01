from pathlib import Path


def parse_to_markdown(input_path: Path, output_path: Path) -> str:
    input_path = Path(input_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    ext = input_path.suffix.lower()

    if ext == ".md":
        md = input_path.read_text(encoding="utf-8")
        output_path.write_text(md, encoding="utf-8")
        return md

    if ext == ".txt":
        text = input_path.read_text(encoding="utf-8")
        output_path.write_text(text, encoding="utf-8")
        return text

    if ext == ".pdf":
        md = _pdf_to_markdown(input_path)
    elif ext == ".docx":
        md = _docx_to_markdown(input_path)
    elif ext == ".pptx":
        md = _pptx_to_markdown(input_path)
    else:
        md = ""

    output_path.write_text(md, encoding="utf-8")
    return md


def _pdf_to_markdown(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(f"## Page {i}\n\n{text}")
    return "\n\n".join(pages)


def _docx_to_markdown(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        style = para.style.name or ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("Heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("Heading 3"):
            parts.append(f"### {text}")
        else:
            parts.append(text)
    return "\n\n".join(parts)


def _pptx_to_markdown(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides)
